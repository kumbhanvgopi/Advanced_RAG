# vectors.py - Improved with better design principles

import os
import re
import logging
import hashlib
import time
import uuid
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from functools import wraps

# Document processing imports with fallbacks
try:
    from langchain_community.document_loaders import PyPDFLoader, UnstructuredPDFLoader
except ImportError:
    from langchain.document_loaders import PyPDFLoader, UnstructuredPDFLoader

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    from langchain.text_splitter import RecursiveCharacterTextSplitter

try:
    from langchain_community.embeddings import HuggingFaceBgeEmbeddings
except ImportError:
    try:
        from langchain_huggingface import HuggingFaceEmbeddings as HuggingFaceBgeEmbeddings
    except ImportError:
        from langchain.embeddings import HuggingFaceEmbeddings as HuggingFaceBgeEmbeddings

try:
    from langchain_community.vectorstores import Qdrant
except ImportError:
    from langchain.vectorstores import Qdrant

from langchain.schema import Document

# Vector database imports
from qdrant_client import QdrantClient
from qdrant_client.http import models
from qdrant_client.models import Distance, VectorParams, PointStruct

# Additional document loaders
try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from dotenv import load_dotenv

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
load_dotenv()

# Cost and Security limits
MAX_DOCUMENT_SIZE = 50 * 1024 * 1024  # 50MB
MAX_CHUNKS_PER_DOCUMENT = 1000
MAX_PROCESSING_TIME = 300  # 5 minutes
ALLOWED_FILE_TYPES = {'.pdf', '.txt', '.docx', '.pptx'}

def retry_on_failure(max_retries=3, delay=1):
    """Decorator for retrying failed operations"""
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            for attempt in range(max_retries):
                try:
                    return func(*args, **kwargs)
                except Exception as e:
                    if attempt == max_retries - 1:
                        raise e
                    logger.warning(f"Attempt {attempt + 1} failed: {e}. Retrying in {delay}s...")
                    time.sleep(delay * (attempt + 1))
            return None
        return wrapper
    return decorator

def validate_file_security(file_path: str) -> Tuple[bool, str]:
    """Security validation for uploaded files"""
    try:
        path = Path(file_path)
        
        # Check file existence
        if not path.exists():
            return False, "File does not exist"
        
        # Check file size
        file_size = path.stat().st_size
        if file_size > MAX_DOCUMENT_SIZE:
            return False, f"File too large: {file_size / (1024*1024):.1f}MB > {MAX_DOCUMENT_SIZE / (1024*1024)}MB"
        
        # Check file extension
        if path.suffix.lower() not in ALLOWED_FILE_TYPES:
            return False, f"File type not allowed: {path.suffix}"
        
        # Check for suspicious file patterns
        if any(pattern in path.name.lower() for pattern in ['../', '..\\', '<script', 'javascript:']):
            return False, "Suspicious file name detected"
        
        return True, "File validation passed"
        
    except Exception as e:
        return False, f"Validation error: {str(e)}"

class EmbeddingsManager:
    """Document-focused embeddings manager applying all 5 design principles."""
    
    def __init__(
        self,
        model_name: str = "BAAI/bge-small-en",
        device: str = "cpu",
        encode_kwargs: dict = None,
        qdrant_url: str = None,
        collection_name: str = None,  # Make it required or session-based
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        max_chunks: int = MAX_CHUNKS_PER_DOCUMENT
    ):
        """Initialize with improved validation and defaults"""
        
        # Validate and sanitize inputs
        self.model_name = self._sanitize_model_name(model_name)
        self.device = device if device in ["cpu", "cuda"] else "cpu"
        self.encode_kwargs = encode_kwargs or {"normalize_embeddings": True}
        self.qdrant_url = qdrant_url or os.getenv('QDRANT_URL')
        
        # IMPORTANT: Use session-specific collection or default
        if collection_name:
            self.collection_name = self._sanitize_collection_name(collection_name)
        else:
            # Generate unique collection name if none provided
            unique_id = str(uuid.uuid4())[:8]
            self.collection_name = f"doc_collection_{unique_id}"
            
        self.api_key = os.getenv('QDRANT_API_KEY')
        
        # Validate and set processing parameters
        self.chunk_size = max(100, min(chunk_size, 2000))
        self.chunk_overlap = max(0, min(chunk_overlap, self.chunk_size // 2))
        self.max_chunks = max_chunks
        
        # Processing statistics
        self.stats = {
            'documents_processed': 0,
            'total_chunks_created': 0,
            'processing_time': 0,
            'errors_encountered': 0
        }
        
        self._validate_config()
        self._initialize_components()

    def _sanitize_model_name(self, model_name: str) -> str:
        """Sanitize model name for security"""
        # Allow only alphanumeric, hyphens, underscores, and slashes
        sanitized = re.sub(r'[^a-zA-Z0-9\-_/]', '', model_name)
        return sanitized if sanitized else "BAAI/bge-small-en"

    def _sanitize_collection_name(self, collection_name: str) -> str:
        """Sanitize collection name for Qdrant"""
        # Replace invalid characters with underscores
        sanitized = re.sub(r'[^a-zA-Z0-9_\-]', '_', collection_name)
        return sanitized[:64]  # Limit length

    def _validate_config(self) -> None:
        """Validate configuration with improved error messages"""
        if not self.qdrant_url:
            raise ValueError("QDRANT_URL must be provided in environment variables")
        
        if not self.api_key:
            logger.warning("QDRANT_API_KEY not found - attempting anonymous connection")
        
        # Validate URL format
        if not (self.qdrant_url.startswith('http://') or self.qdrant_url.startswith('https://')):
            raise ValueError("QDRANT_URL must be a valid HTTP/HTTPS URL")

    @retry_on_failure(max_retries=3)
    def _initialize_components(self) -> None:
        """Initialize all components with retry logic"""
        self._initialize_qdrant_client()
        self._initialize_embeddings()
        self._initialize_text_splitter()
        self._initialize_collection()

    def _initialize_qdrant_client(self) -> None:
        """Initialize Qdrant client with improved error handling"""
        try:
            self.qdrant_client = QdrantClient(
                url=self.qdrant_url,
                api_key=self.api_key,
                timeout=30,
                prefer_grpc=False
            )
            
            # Test connection with timeout
            collections = self.qdrant_client.get_collections()
            logger.info(f"Connected to Qdrant: {len(collections.collections)} collections found")
            
        except Exception as e:
            logger.error(f"Qdrant connection failed: {e}")
            raise ConnectionError(f"Cannot connect to Qdrant at {self.qdrant_url}: {e}")

    def _initialize_embeddings(self) -> None:
        """Initialize embedding model with fallback"""
        try:
            self.embeddings = HuggingFaceBgeEmbeddings(
                model_name=self.model_name,
                model_kwargs={"device": self.device},
                encode_kwargs=self.encode_kwargs,
            )
            logger.info(f"Initialized embedding model: {self.model_name}")
            
        except Exception as e:
            logger.warning(f"Primary model failed: {e}. Trying fallback...")
            try:
                self.embeddings = HuggingFaceBgeEmbeddings(
                    model_name="sentence-transformers/all-MiniLM-L6-v2",
                    model_kwargs={"device": "cpu"},
                    encode_kwargs={"normalize_embeddings": True},
                )
                logger.info("Using fallback model: all-MiniLM-L6-v2")
            except Exception as fallback_error:
                raise RuntimeError(f"All embedding models failed: {fallback_error}")

    def _initialize_text_splitter(self) -> None:
        """Initialize text splitter with improved separators"""
        self.text_splitter = RecursiveCharacterTextSplitter(
            separators=["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " ", ""],
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            length_function=len,
            keep_separator=True
        )
        logger.info(f"Text splitter initialized: {self.chunk_size} chars, {self.chunk_overlap} overlap")

    def _initialize_collection(self) -> None:
        """Initialize or verify collection with improved error handling"""
        try:
            collections = self.qdrant_client.get_collections().collections
            collection_exists = any(col.name == self.collection_name for col in collections)
            
            if not collection_exists:
                vector_config = VectorParams(
                    size=384,  # BGE-small-en size
                    distance=Distance.COSINE
                )
                self.qdrant_client.create_collection(
                    collection_name=self.collection_name,
                    vectors_config=vector_config
                )
                logger.info(f"Created collection: {self.collection_name}")
            else:
                logger.info(f"Using existing collection: {self.collection_name}")
                
        except Exception as e:
            raise ConnectionError(f"Collection initialization failed: {e}")

    def _load_pdf_document(self, file_path: str) -> List[Document]:
        """Load PDF with multiple fallback methods"""
        loaders = [
            (PyPDFLoader, "PyPDFLoader"),
            (UnstructuredPDFLoader, "UnstructuredPDFLoader")
        ]
        
        for loader_class, loader_name in loaders:
            try:
                loader = loader_class(str(file_path))
                documents = loader.load()
                
                if documents and any(doc.page_content.strip() for doc in documents):
                    logger.info(f"PDF loaded with {loader_name}: {len(documents)} pages")
                    return documents
                    
            except Exception as e:
                logger.warning(f"{loader_name} failed: {e}")
                continue
        
        raise ValueError("All PDF loaders failed - file may be corrupted or encrypted")

    def _load_docx_document(self, file_path: str) -> List[Document]:
        """Load DOCX with improved error handling"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx required: pip install python-docx")
        
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(str(file_path))
            
            content_parts = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                text = para.text.strip()
                if text and len(text) > 10:  # Filter very short paragraphs
                    content_parts.append(text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        content_parts.append(row_text)
            
            if not content_parts:
                raise ValueError("No readable content found in DOCX")
            
            combined_text = "\n\n".join(content_parts)
            return [Document(
                page_content=combined_text,
                metadata={"source": str(file_path), "page": 1, "doc_type": "docx"}
            )]
            
        except Exception as e:
            raise ValueError(f"DOCX processing error: {e}")

    def _load_pptx_document(self, file_path: str) -> List[Document]:
        """Load PPTX with slide-by-slide processing"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx required: pip install python-pptx")
        
        try:
            from pptx import Presentation
            ppt = Presentation(str(file_path))
            documents = []
            
            for slide_num, slide in enumerate(ppt.slides, 1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    combined_text = "\n\n".join(slide_text)
                    if len(combined_text.strip()) > 20:  # Filter very short slides
                        documents.append(Document(
                            page_content=combined_text,
                            metadata={
                                "source": str(file_path),
                                "page": slide_num,
                                "slide_number": slide_num,
                                "doc_type": "pptx"
                            }
                        ))
            
            if not documents:
                raise ValueError("No readable slides found in PPTX")
            
            logger.info(f"PPTX loaded: {len(documents)} slides")
            return documents
            
        except Exception as e:
            raise ValueError(f"PPTX processing error: {e}")

    def _load_txt_document(self, file_path: str) -> List[Document]:
        """Load text file with encoding detection"""
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as file:
                    content = file.read()
                
                if not content.strip():
                    raise ValueError("Text file is empty")
                
                # Basic content validation
                if len(content) < 50:
                    raise ValueError("Text file too short (minimum 50 characters)")
                
                return [Document(
                    page_content=content,
                    metadata={"source": str(file_path), "page": 1, "doc_type": "txt", "encoding": encoding}
                )]
                
            except UnicodeDecodeError:
                continue
            except Exception as e:
                raise ValueError(f"TXT processing error: {e}")
        
        raise ValueError("Could not decode text file with any common encoding")

    def _load_document(self, file_path: str) -> List[Document]:
        """Load document with security validation and type detection"""
        
        # Security validation
        is_valid, message = validate_file_security(file_path)
        if not is_valid:
            raise SecurityError(f"File validation failed: {message}")
        
        path = Path(file_path)
        extension = path.suffix.lower()
        
        # Document type mapping
        loaders = {
            ".pdf": self._load_pdf_document,
            ".docx": self._load_docx_document,
            ".pptx": self._load_pptx_document,
            ".txt": self._load_txt_document
        }
        
        if extension not in loaders:
            raise ValueError(f"Unsupported file type: {extension}")
        
        try:
            start_time = time.time()
            documents = loaders[extension](str(file_path))
            processing_time = time.time() - start_time
            
            if processing_time > MAX_PROCESSING_TIME:
                raise TimeoutError(f"Document processing timeout: {processing_time:.1f}s > {MAX_PROCESSING_TIME}s")
            
            # Filter empty documents
            documents = [doc for doc in documents if doc.page_content.strip()]
            
            if not documents:
                raise ValueError("No valid content extracted from document")
            
            logger.info(f"Document loaded: {len(documents)} sections in {processing_time:.2f}s")
            return documents
            
        except Exception as e:
            self.stats['errors_encountered'] += 1
            raise ValueError(f"Document loading failed: {e}")

    def _clean_and_validate_text(self, text: str) -> str:
        """Enhanced text cleaning with security measures"""
        if not text or not text.strip():
            return ""
        
        try:
            # Remove potential security threats
            text = re.sub(r'<script[^>]*>.*?</script>', '[SCRIPT_REMOVED]', text, flags=re.IGNORECASE | re.DOTALL)
            text = re.sub(r'javascript:', '[JS_REMOVED]', text, flags=re.IGNORECASE)
            
            # Basic PII detection and masking
            text = re.sub(r'\b\d{3}[-.\s]?\d{2}[-.\s]?\d{4}\b', '[SSN]', text)
            text = re.sub(r'\b\d{4}[-.\s]?\d{4}[-.\s]?\d{4}[-.\s]?\d{4}\b', '[CARD]', text)
            text = re.sub(r'\b\d{10,11}\b', '[PHONE]', text)
            
            # Clean excessive whitespace
            text = re.sub(r'\s+', ' ', text)
            text = text.strip()
            
            return text
            
        except Exception as e:
            logger.warning(f"Text cleaning failed: {e}")
            return text

    def _enhance_metadata(self, splits: List[Document], file_path: str) -> List[Document]:
        """Enhanced metadata with hash and validation"""
        path = Path(file_path)
        
        # Calculate file hash for uniqueness
        with open(file_path, 'rb') as f:
            file_hash = hashlib.md5(f.read()).hexdigest()
        
        enhanced_splits = []
        for i, split in enumerate(splits):
            # Skip very short chunks
            if len(split.page_content.strip()) < 100:
                continue
            
            # Clean content
            clean_content = self._clean_and_validate_text(split.page_content)
            if not clean_content:
                continue
            
            # Enhanced metadata
            split.page_content = clean_content
            split.metadata.update({
                "file_name": path.name,
                "file_type": path.suffix.lower(),
                "file_hash": file_hash,
                "source": str(path.absolute()),
                "page": split.metadata.get('page', (i // 3 + 1)),
                "chunk_id": i,
                "chunk_size": len(clean_content),
                "word_count": len(clean_content.split()),
                "processing_timestamp": time.time()
            })
            
            enhanced_splits.append(split)
        
        return enhanced_splits

    @retry_on_failure(max_retries=2)
    def create_embeddings(self, file_path: str) -> str:
        """Create embeddings with improved error handling and cost controls"""
        start_time = time.time()
        
        try:
            logger.info(f"Starting embedding creation for: {Path(file_path).name}")
            
            # Load and validate document
            docs = self._load_document(file_path)
            
            # Split into chunks with size limits
            splits = self.text_splitter.split_documents(docs)
            
            # Apply cost controls
            if len(splits) > self.max_chunks:
                logger.warning(f"Too many chunks ({len(splits)}), limiting to {self.max_chunks}")
                splits = splits[:self.max_chunks]
            
            # Enhance metadata and clean content
            splits = self._enhance_metadata(splits, file_path)
            
            if not splits:
                raise ValueError("No valid chunks created after processing")
            
            logger.info(f"Created {len(splits)} valid chunks")
            
            # Create embeddings with fallback methods
            success = False
            methods = [
                ("from_documents", self._create_embeddings_method1),
                ("manual_upsert", self._create_embeddings_method2),
                ("batch_texts", self._create_embeddings_method3)
            ]
            
            for method_name, method_func in methods:
                try:
                    method_func(splits)
                    logger.info(f"Embeddings created using {method_name}")
                    success = True
                    break
                except Exception as e:
                    logger.warning(f"{method_name} failed: {e}")
                    continue
            
            if not success:
                raise ConnectionError("All embedding creation methods failed")
            
            # Update statistics
            processing_time = time.time() - start_time
            self.stats.update({
                'documents_processed': self.stats['documents_processed'] + 1,
                'total_chunks_created': self.stats['total_chunks_created'] + len(splits),
                'processing_time': self.stats['processing_time'] + processing_time
            })
            
            return self._generate_success_message(file_path, splits, processing_time)
            
        except Exception as e:
            self.stats['errors_encountered'] += 1
            logger.error(f"Embedding creation failed: {e}")
            raise Exception(f"Embedding creation failed: {e}")

    def _create_embeddings_method1(self, splits: List[Document]) -> None:
        """Method 1: Standard from_documents"""
        Qdrant.from_documents(
            documents=splits,
            embedding=self.embeddings,
            url=self.qdrant_url,
            api_key=self.api_key,
            collection_name=self.collection_name,
            prefer_grpc=False,
        )

    def _create_embeddings_method2(self, splits: List[Document]) -> None:
        """Method 2: Manual upsert"""
        texts = [doc.page_content for doc in splits]
        embeddings_list = self.embeddings.embed_documents(texts)
        
        points = []
        for i, (doc, embedding) in enumerate(zip(splits, embeddings_list)):
            point = PointStruct(
                id=f"{hash(doc.page_content)}_{i}",
                vector=embedding,
                payload={
                    "page_content": doc.page_content,
                    **doc.metadata
                }
            )
            points.append(point)
        
        self.qdrant_client.upsert(
            collection_name=self.collection_name,
            points=points
        )

    def _create_embeddings_method3(self, splits: List[Document]) -> None:
        """Method 3: Batch texts"""
        texts = [doc.page_content for doc in splits]
        metadatas = [doc.metadata for doc in splits]
        
        Qdrant.from_texts(
            texts=texts,
            embedding=self.embeddings,
            metadatas=metadatas,
            url=self.qdrant_url,
            api_key=self.api_key,
            collection_name=self.collection_name,
            prefer_grpc=False,
        )

    def _generate_success_message(self, file_path: str, splits: List[Document], processing_time: float) -> str:
        """Generate detailed success message"""
        file_name = Path(file_path).name
        total_chunks = len(splits)
        avg_chunk_size = sum(len(split.page_content) for split in splits) / len(splits)
        total_words = sum(split.metadata.get('word_count', 0) for split in splits)
        
        return (
            f"✅ Successfully processed '{file_name}'!\n\n"
            f"📊 Processing Summary:\n"
            f"• Created {total_chunks} text chunks\n"
            f"• Average chunk size: {avg_chunk_size:.0f} characters\n"
            f"• Total words processed: {total_words}\n"
            f"• Processing time: {processing_time:.2f} seconds\n"
            f"• Collection: {self.collection_name}\n"
            f"• Security validation: Passed"
        )

    def get_stats(self) -> Dict[str, Any]:
        """Get processing statistics"""
        return {
            **self.stats,
            "collection_name": self.collection_name,
            "max_chunks_limit": self.max_chunks,
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap
        }

    def clear_collection(self) -> None:
        """Safely clear collection"""
        try:
            self.qdrant_client.delete_collection(self.collection_name)
            self._initialize_collection()
            logger.info(f"Collection cleared: {self.collection_name}")
        except Exception as e:
            logger.error(f"Failed to clear collection: {e}")
            raise

class SecurityError(Exception):
    """Custom exception for security-related errors"""
    pass